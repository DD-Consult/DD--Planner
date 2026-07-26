"""
Iteration 42 — Branded cover slide/page for PPT and PDF exports.

Verifies:
- PPTX project export: first slide is a branded cover (navy full-bleed rect,
  'PROJECT STATUS REPORT' eyebrow, project name title, 'For {client_name}',
  period + report date, 'Prepared by DD Consulting', 'CONFIDENTIAL').
- PPTX WBS export: first slide cover with 'WORK BREAKDOWN STRUCTURE' subtitle.
- PDF export: still 16:9; total pages > 1 (cover + content). The print URL JS
  bundle includes data-export-section="cover".
- _fetch_cover_meta returns fallback dict for invalid project id (no crash).
- Section slides no longer contain "Slide N" text in header.
"""
import io
import os
import re
import sys
import asyncio
import pytest
import requests
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Emu

sys.path.insert(0, "/app/backend")


def _load_backend_url():
    v = os.environ.get("REACT_APP_BACKEND_URL")
    if not v:
        env = "/app/frontend/.env"
        if os.path.exists(env):
            for line in open(env):
                if line.startswith("REACT_APP_BACKEND_URL="):
                    v = line.split("=", 1)[1].strip().strip('"').strip("'")
                    break
    return v.rstrip("/")


BASE_URL = _load_backend_url()
ADMIN = {"email": "admin@test.com", "password": "admin123"}


@pytest.fixture(scope="module")
def admin_headers():
    r = requests.post(
        f"{BASE_URL}/api/auth/login",
        data={"username": ADMIN["email"], "password": ADMIN["password"]},
        headers={"Content-Type": "application/x-www-form-urlencoded"},
        timeout=30,
    )
    assert r.status_code == 200, f"login failed: {r.status_code} {r.text[:200]}"
    tok = r.json().get("access_token") or r.json().get("token")
    return {"Authorization": f"Bearer {tok}"}


@pytest.fixture(scope="module")
def website_project(admin_headers):
    r = requests.get(f"{BASE_URL}/api/projects", headers=admin_headers, timeout=30)
    assert r.status_code == 200
    data = r.json()
    projects = data if isinstance(data, list) else data.get("projects") or data.get("items") or []
    proj = next((p for p in projects if "website" in (p.get("name") or "").lower()), None) or (projects[0] if projects else None)
    assert proj, "no projects available"
    return proj


def _all_text_shapes(slide):
    """Concatenate all text-frame text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        texts.append(run.text)
                if not para.runs and para.text:
                    texts.append(para.text)
    return " || ".join(texts)


class TestProjectPptCover:
    def test_first_slide_is_branded_cover(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/ppt",
            headers=admin_headers,
            timeout=300,
        )
        assert r.status_code == 200, f"PPT export failed: {r.status_code} {r.text[:200]}"
        prs = Presentation(io.BytesIO(r.content))
        n = len(prs.slides)
        print(f"[project ppt] total slides = {n}")
        # 1 cover + N section slides
        assert n >= 2, f"expected cover + content, got {n} slides"

        cover = prs.slides[0]
        text = _all_text_shapes(cover)
        print(f"[project ppt cover text] {text[:500]}")

        # Required cover elements
        assert "PROJECT STATUS REPORT" in text.upper(), "eyebrow missing"
        assert "DD CONSULTING" in text.upper() or "Prepared by DD Consulting" in text, "brand footer missing"
        assert "CONFIDENTIAL" in text.upper(), "CONFIDENTIAL tag missing"
        # Project name (from seeded 'Website Redesign' or first project)
        proj_name = (website_project.get("name") or "").strip()
        if proj_name:
            assert proj_name.split()[0].lower() in text.lower(), f"project name '{proj_name}' not in cover"

        # Full-bleed navy rect: find a rectangle roughly the size of the slide
        sw = prs.slide_width
        sh = prs.slide_height
        found_bg = False
        for shape in cover.shapes:
            # shape_type 1 = AUTO_SHAPE / RECTANGLE
            try:
                if shape.width >= sw * 0.99 and shape.height >= sh * 0.99:
                    found_bg = True
                    break
            except Exception:
                pass
        assert found_bg, "full-bleed background rectangle not found on cover"

    def test_client_name_and_period_on_cover(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/ppt",
            headers=admin_headers,
            timeout=300,
        )
        assert r.status_code == 200
        prs = Presentation(io.BytesIO(r.content))
        text = _all_text_shapes(prs.slides[0])
        client = website_project.get("client_name")
        if client:
            assert client in text, f"client name '{client}' missing on cover; got: {text[:400]}"
            assert "For " in text, "'For {client}' phrasing missing"
        # Period should look like 'Mon YYYY – Mon YYYY' if start/end dates exist
        if website_project.get("start_date") and website_project.get("end_date"):
            assert re.search(r"[A-Z][a-z]{2}\s+\d{4}\s*[–-]\s*[A-Z][a-z]{2}\s+\d{4}", text), \
                f"period range not found: {text[:400]}"

    def test_section_slides_have_no_slide_number_prefix(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/ppt",
            headers=admin_headers,
            timeout=300,
        )
        prs = Presentation(io.BytesIO(r.content))
        # Skip cover (slide 0), check section slides
        for idx in range(1, len(prs.slides)):
            slide = prs.slides[idx]
            text = _all_text_shapes(slide)
            assert not re.search(r"\bSlide\s+\d+\b", text), \
                f"section slide {idx} still says 'Slide N': {text[:200]}"


class TestWbsPptCover:
    def test_wbs_cover_has_wbs_subtitle(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/wbs/ppt",
            headers=admin_headers,
            timeout=300,
        )
        assert r.status_code == 200
        prs = Presentation(io.BytesIO(r.content))
        assert len(prs.slides) >= 2, "wbs pptx should have cover + content"
        cover_text = _all_text_shapes(prs.slides[0])
        print(f"[wbs cover text] {cover_text[:400]}")
        assert "WORK BREAKDOWN STRUCTURE" in cover_text.upper(), \
            f"WBS subtitle missing: {cover_text[:300]}"
        assert "CONFIDENTIAL" in cover_text.upper()


class TestPdfCover:
    def test_pdf_has_cover_page(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/pdf",
            headers=admin_headers,
            timeout=300,
        )
        assert r.status_code == 200
        assert r.content[:5] == b"%PDF-"
        reader = PdfReader(io.BytesIO(r.content))
        n = len(reader.pages)
        print(f"[pdf pages] {n}")
        assert n >= 2, f"expected cover + content pages, got {n}"
        # Still 16:9
        box = reader.pages[0].mediabox
        w, h = float(box.width), float(box.height)
        assert 940 <= w <= 980 and 520 <= h <= 560, f"cover page not 16:9: {w}x{h}"

    def test_print_html_contains_cover_marker(self, admin_headers, website_project):
        """The cover div lives in the JS bundle since ProjectReport is SPA."""
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(f"{BASE_URL}/projects/{pid}/report", timeout=30)
        assert r.status_code == 200
        html = r.text
        js_srcs = re.findall(r'src="([^"]+\.js)"', html)
        found = False
        for src in js_srcs[:20]:
            url = src if src.startswith("http") else BASE_URL + src
            try:
                jr = requests.get(url, timeout=30)
                if jr.status_code != 200:
                    continue
                body = jr.text
                if ('data-export-section="cover"' in body
                    or "data-export-section='cover'" in body
                    or 'data-export-section=\\"cover\\"' in body
                    or 'PROJECT STATUS REPORT' in body):
                    found = True
                    break
            except Exception:
                continue
        assert found, "cover marker not found in JS bundles"


class TestFetchCoverMetaResilience:
    def test_invalid_project_id_returns_fallback(self):
        from services.exports.ppt_export import _fetch_cover_meta

        async def _run():
            return await _fetch_cover_meta("nonexistent_id_xyz", subtitle="TEST SUBTITLE")

        result = asyncio.run(_run())
        assert isinstance(result, dict), "should return a dict, not crash"
        assert result.get("project_name") == "Project Report"
        assert result.get("subtitle") == "TEST SUBTITLE"
        assert "report_date" in result

    def test_valid_project_id_returns_full_meta(self, admin_headers, website_project):
        from services.exports.ppt_export import _fetch_cover_meta

        pid = website_project.get("id") or website_project.get("_id")

        async def _run():
            return await _fetch_cover_meta(pid, subtitle="PROJECT STATUS REPORT")

        result = asyncio.run(_run())
        assert isinstance(result, dict)
        assert result.get("subtitle") == "PROJECT STATUS REPORT"
        assert result.get("project_name")
        print(f"[cover meta] {result}")
