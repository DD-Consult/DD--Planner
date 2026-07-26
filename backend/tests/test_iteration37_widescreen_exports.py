"""
Iteration 37 — 16:9 widescreen exports (PDF + PPTX) and WBS-in-report

Verifies:
- PDF exports use 16:9 (~960 x 540 pt) MediaBox, not A4
- PPTX exports use 13.333in x 7.5in slide dims and produce multiple section slides
  with picture shapes preserving original aspect ratios (not stretched to fill).
- WBS PDF/PPT variants use 16:9
- Project Report HTML (print URL) contains the WBS section container
- Permission: resource user should not access admin export endpoints (or gets
  their own subset). We record actual behavior.
"""
import io
import os
import re
import time
import pytest
import requests

from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Emu

def _load_backend_url():
    v = os.environ.get("REACT_APP_BACKEND_URL")
    if not v:
        env = "/app/frontend/.env"
        if os.path.exists(env):
            for line in open(env):
                if line.startswith("REACT_APP_BACKEND_URL="):
                    v = line.split("=", 1)[1].strip().strip('"').strip("'")
                    break
    if not v:
        raise RuntimeError("REACT_APP_BACKEND_URL not set")
    return v.rstrip("/")


BASE_URL = _load_backend_url()

ADMIN = {"email": "admin@test.com", "password": "admin123"}
RESOURCE = {"email": "riley@test.com", "password": "riley123"}


# -------------------- fixtures --------------------
@pytest.fixture(scope="module")
def admin_token():
    r = requests.post(
        f"{BASE_URL}/api/auth/login",
        data={"username": ADMIN["email"], "password": ADMIN["password"]},
        headers={"Content-Type": "application/x-www-form-urlencoded"},
        timeout=30,
    )
    assert r.status_code == 200, f"admin login failed: {r.status_code} {r.text[:200]}"
    return r.json().get("access_token") or r.json().get("token")


@pytest.fixture(scope="module")
def resource_token():
    r = requests.post(
        f"{BASE_URL}/api/auth/login",
        data={"username": RESOURCE["email"], "password": RESOURCE["password"]},
        headers={"Content-Type": "application/x-www-form-urlencoded"},
        timeout=30,
    )
    if r.status_code != 200:
        pytest.skip(f"resource login not available: {r.status_code}")
    return r.json().get("access_token") or r.json().get("token")


@pytest.fixture(scope="module")
def admin_headers(admin_token):
    return {"Authorization": f"Bearer {admin_token}"}


@pytest.fixture(scope="module")
def website_project(admin_headers):
    r = requests.get(f"{BASE_URL}/api/projects", headers=admin_headers, timeout=30)
    assert r.status_code == 200, f"projects list failed: {r.status_code}"
    data = r.json()
    projects = data if isinstance(data, list) else data.get("projects") or data.get("items") or []
    # prefer "Website Redesign" (seeded project mentioned in context)
    proj = next((p for p in projects if "website" in (p.get("name") or "").lower()), None) or (projects[0] if projects else None)
    assert proj, "No projects available for testing"
    return proj


# -------------------- helpers --------------------
def _mediabox_size(pdf_bytes):
    reader = PdfReader(io.BytesIO(pdf_bytes))
    assert len(reader.pages) >= 1
    box = reader.pages[0].mediabox
    w = float(box.width)
    h = float(box.height)
    return w, h


# -------------------- PDF tests --------------------
class TestProjectPdfWidescreen:
    def test_pdf_export_is_16_9(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/pdf",
            headers=admin_headers,
            timeout=180,
        )
        assert r.status_code == 200, f"PDF export failed: {r.status_code} {r.text[:200]}"
        assert r.headers.get("content-type", "").startswith("application/pdf"), \
            f"Wrong content-type: {r.headers.get('content-type')}"
        assert r.content[:5] == b"%PDF-", "PDF magic bytes missing"
        assert len(r.content) > 10_000, f"PDF too small: {len(r.content)}"

        w, h = _mediabox_size(r.content)
        # 13.333in * 72 = 960, 7.5in * 72 = 540
        aspect = w / h
        print(f"[project pdf] mediabox = {w:.1f} x {h:.1f} pt, aspect={aspect:.3f}")
        # Not A4 landscape (~842 x 595)
        assert not (820 < w < 860 and 580 < h < 610), \
            f"MediaBox looks like A4 landscape: {w}x{h}"
        # ~ 960 x 540 (allow reasonable tolerance)
        assert 940 <= w <= 980, f"Width not ~960pt: {w}"
        assert 520 <= h <= 560, f"Height not ~540pt: {h}"
        assert 1.75 <= aspect <= 1.80, f"Aspect not 16:9: {aspect}"

    def test_wbs_pdf_export_is_16_9(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/wbs/pdf",
            headers=admin_headers,
            timeout=180,
        )
        assert r.status_code == 200, f"WBS PDF failed: {r.status_code} {r.text[:200]}"
        assert r.content[:5] == b"%PDF-"
        w, h = _mediabox_size(r.content)
        aspect = w / h
        print(f"[wbs pdf] mediabox = {w:.1f} x {h:.1f} pt, aspect={aspect:.3f}")
        assert 940 <= w <= 980, f"Width not ~960pt: {w}"
        assert 520 <= h <= 560, f"Height not ~540pt: {h}"
        assert 1.75 <= aspect <= 1.80, f"Aspect not 16:9: {aspect}"


# -------------------- PPTX tests --------------------
class TestProjectPptWidescreen:
    def test_ppt_export_is_16_9_multi_slide_preserved_aspect(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/ppt",
            headers=admin_headers,
            timeout=240,
        )
        assert r.status_code == 200, f"PPT export failed: {r.status_code} {r.text[:200]}"
        ct = r.headers.get("content-type", "")
        assert "presentationml" in ct, f"Wrong content-type: {ct}"
        assert len(r.content) > 50_000, f"PPTX too small: {len(r.content)}"

        prs = Presentation(io.BytesIO(r.content))
        # Slide size 13.333in × 7.5in — python-pptx uses EMU. 1in = 914400 EMU
        sw = prs.slide_width
        sh = prs.slide_height
        print(f"[ppt] slide_size EMU = {sw} x {sh}  ({sw/914400:.3f}in x {sh/914400:.3f}in)")
        assert abs(sw - Emu(13.333 * 914400)) < 20_000, f"slide_width not 13.333in: {sw}"
        assert abs(sh - Emu(7.5 * 914400)) < 20_000, f"slide_height not 7.5in: {sh}"

        # Multi-section slides expected: header + at least 3 section slides
        n = len(prs.slides)
        print(f"[ppt] slide count = {n}")
        assert n >= 3, f"Expected multi-slide PPTX, got {n}"

        # For each picture shape, ensure it is NOT stretched to a fixed 12.933 x 6.7 in area.
        # The old (broken) behavior scaled every image to exactly 12.933in x 6.7in
        # regardless of source aspect ratio.
        BAD_W = int(12.933 * 914400)
        BAD_H = int(6.7 * 914400)
        pic_count = 0
        stretched_pics = 0
        for idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    pic_count += 1
                    w_emu = shape.width
                    h_emu = shape.height
                    # tolerance +/- 0.02 in
                    tol = int(0.02 * 914400)
                    if abs(w_emu - BAD_W) < tol and abs(h_emu - BAD_H) < tol:
                        stretched_pics += 1
                    print(f"  slide {idx} picture: {w_emu/914400:.3f}in x {h_emu/914400:.3f}in")
        assert pic_count >= 1, "No picture shapes found in PPTX"
        # If more than one picture, at least one should have varied dimensions (not all stretched)
        if pic_count >= 2:
            assert stretched_pics < pic_count, (
                f"All {pic_count} pictures have the fixed 12.933x6.7in size — aspect ratio not preserved"
            )

    def test_wbs_ppt_export_is_16_9(self, admin_headers, website_project):
        pid = website_project.get("id") or website_project.get("_id")
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/wbs/ppt",
            headers=admin_headers,
            timeout=240,
        )
        assert r.status_code == 200, f"WBS PPT failed: {r.status_code} {r.text[:200]}"
        prs = Presentation(io.BytesIO(r.content))
        sw = prs.slide_width
        sh = prs.slide_height
        print(f"[wbs ppt] slide_size = {sw/914400:.3f}in x {sh/914400:.3f}in, slides={len(prs.slides)}")
        assert abs(sw - Emu(13.333 * 914400)) < 20_000
        assert abs(sh - Emu(7.5 * 914400)) < 20_000
        assert len(prs.slides) >= 1


# -------------------- Report body contains WBS section --------------------
class TestReportBodyHasWBS:
    def test_report_html_contains_wbs_section(self, admin_headers, website_project):
        """Fetch the front-end print URL for the project report and verify WBS container.

        This uses the same base URL (kubernetes ingress) since front-end and
        back-end share host in this environment. If the front-end route returns
        the SPA shell without SSR, we at least verify the /projects/*/report
        route responds 200 and returns HTML. The actual data-export-section='wbs'
        marker lives in the JS bundle, so we also grep the served JS.
        """
        pid = website_project.get("id") or website_project.get("_id")
        # Try the front-end URL (frontend uses same base). Path handled by SPA.
        r = requests.get(f"{BASE_URL}/projects/{pid}/report", timeout=30)
        # SPA may return 200 with index.html
        assert r.status_code == 200, f"report page unreachable: {r.status_code}"
        html = r.text
        assert "<html" in html.lower(), "Not HTML response from report route"

        # Grep the main bundle referenced in the HTML for the marker.
        js_srcs = re.findall(r'src="([^"]+\.js)"', html)
        found_marker = False
        for src in js_srcs[:20]:
            js_url = src if src.startswith("http") else BASE_URL + src
            try:
                jr = requests.get(js_url, timeout=30)
                if jr.status_code == 200 and 'data-export-section="wbs"' in jr.text:
                    found_marker = True
                    break
                # also allow single-quoted / minified variant
                if jr.status_code == 200 and "data-export-section=\\\"wbs\\\"" in jr.text:
                    found_marker = True
                    break
                if jr.status_code == 200 and 'data-export-section=wbs' in jr.text:
                    found_marker = True
                    break
            except Exception:
                continue
        assert found_marker, "WBS section marker not found in any served JS bundle"


# -------------------- Permissions --------------------
class TestExportPermissions:
    def test_resource_pdf_access(self, resource_token, website_project):
        """Record whether the resource role can invoke the PDF export."""
        pid = website_project.get("id") or website_project.get("_id")
        headers = {"Authorization": f"Bearer {resource_token}"}
        r = requests.get(
            f"{BASE_URL}/api/projects/{pid}/export/pdf",
            headers=headers,
            timeout=180,
        )
        # Not asserting a strict permission — just document
        print(f"[resource pdf] status={r.status_code}")
        assert r.status_code in (200, 401, 403, 404), f"Unexpected status: {r.status_code}"
