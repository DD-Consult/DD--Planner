"""
PowerPoint export service using Playwright screenshots + python-pptx composition.

Step 9 update: cover slide branding (colors, company name) is now read from
the current tenant's `branding` and `name` fields rather than hardcoded to
"DD Consulting". Callers can pass `branding` in the `cover` dict to override.
When branding is missing, the historical DD Navy/Gold palette is used as a
safe default so this file remains backward compatible.
"""
import logging
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .renderer import render_screenshots

logger = logging.getLogger(__name__)

# Default (DD Consulting) brand palette — used when a tenant has no branding.
DD_NAVY = RGBColor(0x1B, 0x2A, 0x47)
DD_GOLD = RGBColor(0xC9, 0xA8, 0x4C)
DD_LIGHT = RGBColor(0xE8, 0xED, 0xF2)


def _hex_to_rgb(hex_str: str, fallback: RGBColor) -> RGBColor:
    """Convert '#RRGGBB' to python-pptx RGBColor. Falls back on any error."""
    try:
        if not hex_str or not hex_str.startswith("#") or len(hex_str) != 7:
            return fallback
        return RGBColor(int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
    except Exception:
        return fallback


def _resolve_brand_palette(branding: dict) -> tuple[RGBColor, RGBColor, RGBColor]:
    """Return (primary, accent, light) RGBColors for the cover slide.
    
    `branding` is expected to look like:
        {"primary_color": "#1B2A47", "accent_color": "#C9A84C", "logo_url": "..."}
    Missing/invalid values fall back to the DD palette.
    """
    branding = branding or {}
    primary = _hex_to_rgb(branding.get("primary_color"), DD_NAVY)
    accent = _hex_to_rgb(branding.get("accent_color"), DD_GOLD)
    # No user-configurable "light" yet — derive from primary by lightening,
    # but for now keep the DD light as safe default.
    light = DD_LIGHT
    return primary, accent, light


def _add_cover_slide(prs: Presentation, cover: dict) -> None:
    """Prepend a branded cover slide. `cover` may contain:
       - project_name (str)
       - client_name (str)
       - period (str, e.g. "Q1 2026" or "Jan – Mar 2026")
       - report_date (str, e.g. "March 28, 2026")
       - subtitle (str, optional, e.g. "Weekly Status Report")
       - branding (dict, optional): {primary_color, accent_color, logo_url}
       - workspace_name (str, optional): shown in footer as "Prepared by X"
    """
    from pptx.enum.shapes import MSO_SHAPE

    # Resolve per-tenant brand palette (falls back to DD colors)
    primary, accent, light = _resolve_brand_palette(cover.get("branding") or {})
    # "Prepared by X" — use tenant name if provided, else default to DD
    workspace_name = cover.get("workspace_name") or "DD Consulting"

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full-bleed primary background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()

    # Left-side accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75), Inches(1.2),
        Inches(0.15), Inches(5.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Small "PROJECT REPORT" eyebrow label
    eyebrow = slide.shapes.add_textbox(Inches(1.1), Inches(1.15), Inches(11), Inches(0.4))
    tf = eyebrow.text_frame
    tf.text = (cover.get("subtitle") or "PROJECT REPORT").upper()
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent

    # Project name (huge)
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.75), Inches(11.2), Inches(1.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = cover.get("project_name") or "Project Report"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Client name
    if cover.get("client_name"):
        client_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(11), Inches(0.7))
        tf = client_box.text_frame
        tf.text = f"For {cover['client_name']}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = light

    # Period + Report date
    meta_parts = []
    if cover.get("period"):
        meta_parts.append(cover["period"])
    if cover.get("report_date"):
        meta_parts.append(cover["report_date"])
    if meta_parts:
        meta_box = slide.shapes.add_textbox(Inches(1.1), Inches(4.7), Inches(11), Inches(0.5))
        tf = meta_box.text_frame
        tf.text = "  ·  ".join(meta_parts)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = accent

    # Bottom band — slightly darker shade of primary
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.9),
        prs.slide_width, Inches(0.6),
    )
    band.fill.solid()
    # Darken primary by 15% for the footer band (simple approximation).
    # If we can't compute (fallback), use the historical dark navy.
    try:
        r, g, b = primary[0], primary[1], primary[2]
        band.fill.fore_color.rgb = RGBColor(max(0, r - 12), max(0, g - 15), max(0, b - 23))
    except Exception:
        band.fill.fore_color.rgb = RGBColor(0x0F, 0x1B, 0x30)
    band.line.fill.background()

    brand = slide.shapes.add_textbox(Inches(0.75), Inches(7.02), Inches(6), Inches(0.4))
    tf = brand.text_frame
    tf.text = f"Prepared by {workspace_name}"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = accent

    confidential = slide.shapes.add_textbox(prs.slide_width - Inches(4), Inches(7.02), Inches(3.25), Inches(0.4))
    tf = confidential.text_frame
    tf.text = "CONFIDENTIAL"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = light
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def _compose_pptx(images: list, title: str, cover: dict = None) -> bytes:
    """
    Compose a PowerPoint presentation from a list of screenshot images.
    Each image gets its own 16:9 slide with proper aspect-ratio preservation
    (fit-to-contain rather than stretch-to-fill). When `cover` metadata is
    provided, a branded cover slide is prepended.
    
    Args:
        images: List of PNG image bytes
        title: Presentation title (used in header of each section slide)
        cover: Optional dict with project_name / client_name / period / report_date / subtitle
    
    Returns:
        PPTX bytes
    """
    from PIL import Image as PILImage

    logger.info(f"Composing PPTX with {len(images)} images, title: {title}, cover: {bool(cover)}")

    # Create presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add branded cover as first slide if metadata provided
    if cover:
        _add_cover_slide(prs, cover)

    header_height = Inches(0.55)

    # Resolve per-tenant palette for section slides too (same defaults as cover)
    section_primary, section_accent, _section_light = _resolve_brand_palette((cover or {}).get("branding") or {})
    section_workspace_name = (cover or {}).get("workspace_name") or "DD Consulting"

    # Content area (below header, small margins)
    content_left = Inches(0.2)
    content_top = header_height + Inches(0.1)
    content_max_width = prs.slide_width - Inches(0.4)   # 12.933"
    content_max_height = prs.slide_height - header_height - Inches(0.25)  # ~6.7"

    for idx, image_bytes in enumerate(images):
        logger.info(f"Adding slide {idx + 1} of {len(images)}")

        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            prs.slide_width, header_height
        )
        header.fill.solid()
        header.fill.fore_color.rgb = section_primary
        header.line.fill.background()

        # Title in header
        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.12),
            Inches(10), Inches(0.32)
        )
        text_frame = title_box.text_frame
        text_frame.text = title
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Workspace name on the right (was hardcoded 'DD Consulting')
        logo_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(2.8), Inches(0.14),
            Inches(2.5), Inches(0.3)
        )
        logo_frame = logo_box.text_frame
        logo_frame.text = section_workspace_name
        logo_frame.paragraphs[0].font.size = Pt(13)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = section_accent
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # ── FIT image with preserved aspect ratio ──
        # Read image dimensions to compute correct scale
        try:
            img = PILImage.open(BytesIO(image_bytes))
            img_w_px, img_h_px = img.size
        except Exception as e:
            logger.warning(f"Could not read image dims: {e}; using slide dims")
            img_w_px, img_h_px = (1600, 900)

        # Scale to fit content area, keeping aspect ratio
        max_w = content_max_width
        max_h = content_max_height
        # Determine scale using EMU-ratio (Inches are EMU-based, so we compare ratios)
        img_ratio = img_w_px / img_h_px
        slot_ratio = max_w / max_h

        if img_ratio >= slot_ratio:
            # Width-bound
            new_w = max_w
            new_h = int(max_w / img_ratio)
        else:
            # Height-bound
            new_h = max_h
            new_w = int(max_h * img_ratio)

        # Center the image within the content area
        left = content_left + int((max_w - new_w) / 2)
        top = content_top + int((max_h - new_h) / 2)

        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream,
            left, top,
            width=new_w,
            height=new_h,
        )

        logger.info(f"Slide {idx + 1} added ({new_w} × {new_h} EMU)")

    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_bytes = pptx_stream.getvalue()
    logger.info(f"PPTX generated: {len(pptx_bytes)} bytes")
    return pptx_bytes


async def _fetch_cover_meta(project_id: str, subtitle: str = "PROJECT REPORT") -> dict:
    """Fetch project metadata + tenant branding for the cover slide.
    
    Step 9: also loads the current tenant's `branding` and `name` from
    platform_db so exports show the correct workspace identity. Falls back
    to the historical DD Consulting branding when tenant lookup fails.
    """
    from datetime import datetime
    from bson import ObjectId
    from database import projects_collection

    # --- Tenant branding lookup (best-effort; safe fallback to DD defaults) ---
    branding: dict = {}
    workspace_name = "DD Consulting"
    try:
        from platform_db import tenants_collection as _tenants
        # Prefer the default tenant (which is DD in single-tenant mode). In
        # multi-tenant mode this will resolve to whichever tenant the request
        # was scoped to via the LazyCollection, so we simply read the default.
        # Note: exports run in a background context without a Request object,
        # so tenant-aware routing must be provided by the caller in future.
        # For now, always use the default tenant (DD Consulting).
        _t = await _tenants.find_one({"is_default": True})
        if not _t:
            _t = await _tenants.find_one({})
        if _t:
            workspace_name = _t.get("name") or workspace_name
            branding = _t.get("branding") or {}
    except Exception as _e:
        logger.warning(f"[EXPORT] Could not resolve tenant branding: {_e}; using DD defaults")

    try:
        p = await projects_collection.find_one({"_id": ObjectId(project_id)})
    except Exception:
        p = None
    if not p:
        return {
            "project_name": "Project Report",
            "subtitle": subtitle,
            "report_date": datetime.now().strftime("%B %d, %Y"),
            "workspace_name": workspace_name,
            "branding": branding,
        }

    # Period label based on project dates
    period = None
    s = p.get("start_date")
    e = p.get("end_date")
    if s and e:
        try:
            if isinstance(s, str):
                s = datetime.fromisoformat(s.replace("Z", "+00:00"))
            if isinstance(e, str):
                e = datetime.fromisoformat(e.replace("Z", "+00:00"))
            period = f"{s.strftime('%b %Y')} – {e.strftime('%b %Y')}"
        except Exception:
            period = None

    return {
        "project_name": p.get("name") or "Project Report",
        "client_name": p.get("client_name"),
        "period": period,
        "report_date": datetime.now().strftime("%B %d, %Y"),
        "subtitle": subtitle,
        "workspace_name": workspace_name,
        "branding": branding,
    }


async def build_project_ppt(project_id: str, token: str, frontend_base_url: str) -> bytes:
    """
    Generate PowerPoint for a project report.
    Uses per-section screenshots (data-export-section) so each section becomes
    its own 16:9 slide — avoids the "one giant squished image" problem.
    Prepended with a branded cover slide.
    
    Args:
        project_id: The project ID
        token: JWT token for authentication
        frontend_base_url: Base URL of the frontend (e.g., http://localhost:3000)
    
    Returns:
        PPTX bytes
    """
    logger.info(f"Building project PPT for project_id={project_id}")
    url = f"{frontend_base_url}/print/projects/{project_id}/report?print=1&_t={token}"
    
    # Screenshot each report section separately → one slide per section.
    section_selectors = [
        "#report-header",
        "[data-export-section='summary']",
        "[data-export-section='timeline']",
        "[data-export-section='overview']",
        "[data-export-section='budget']",
        "[data-print-section='risks']",
        "[data-export-section='wbs']",
    ]

    pngs = await render_screenshots(
        url,
        viewport={'width': 1600, 'height': 900},
        selectors=section_selectors,
    )

    if not pngs:
        logger.warning("No section screenshots captured; falling back to full page")
        pngs = await render_screenshots(url, viewport={'width': 1600, 'height': 900})

    if not pngs:
        raise ValueError("Failed to generate project report screenshots")

    cover_meta = await _fetch_cover_meta(project_id, subtitle="PROJECT STATUS REPORT")
    pptx_bytes = _compose_pptx(pngs, title=cover_meta.get("project_name") or "Project Report", cover=cover_meta)
    
    logger.info(f"Project PPT generated: {len(pptx_bytes)} bytes, {len(pngs)} content slide(s) + 1 cover")
    return pptx_bytes


async def build_wbs_ppt(project_id: str, token: str, frontend_base_url: str) -> bytes:
    """
    Generate PowerPoint for a WBS (Work Breakdown Structure).
    Prepended with a branded cover slide.
    """
    logger.info(f"Building WBS PPT for project_id={project_id}")
    url = f"{frontend_base_url}/print/projects/{project_id}/report?print=1&view=wbs&_t={token}"
    
    pngs = await render_screenshots(
        url,
        viewport={'width': 1600, 'height': 900},
        selectors=["#report-header", "[data-export-section='wbs']"],
    )
    if not pngs:
        pngs = await render_screenshots(url, viewport={'width': 1600, 'height': 900})
    
    if not pngs:
        logger.warning("No screenshots captured, WBS report may be empty")
        raise ValueError("Failed to generate WBS report screenshots")
    
    cover_meta = await _fetch_cover_meta(project_id, subtitle="WORK BREAKDOWN STRUCTURE")
    pptx_bytes = _compose_pptx(pngs, title="Work Breakdown Structure", cover=cover_meta)
    
    logger.info(f"WBS PPT generated: {len(pptx_bytes)} bytes, {len(pngs)} content slide(s) + 1 cover")
    return pptx_bytes
